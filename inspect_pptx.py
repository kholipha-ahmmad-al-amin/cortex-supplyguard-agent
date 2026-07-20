import os
from pptx import Presentation

pptx_path = r"D:\all my projects\coco cli compitition\Prototype Submission Template _ Cortex Code CLI Hackathon.pptx"

if not os.path.exists(pptx_path):
    print("File not found:", pptx_path)
else:
    prs = Presentation(pptx_path)
    print(f"Total Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n--- SLIDE {i} ---")
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                print(f"Shape {j} text:\n{shape.text.strip()}")
